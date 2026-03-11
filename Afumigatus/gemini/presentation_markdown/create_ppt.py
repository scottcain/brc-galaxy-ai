from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

def add_slide(prs, title, content_lines, screenshot_placeholder=False):
    layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, line in enumerate(content_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        
        # Handle indentations
        indent_level = 0
        while line.startswith(' '):
            line = line[1:]
            indent_level += 1
            
        p.text = line.strip('* ')
        p.level = indent_level // 2 if indent_level > 0 else 0
        
    if screenshot_placeholder:
        # Add a placeholder box for screenshots
        left = Inches(5.5)
        top = Inches(2.5)
        width = Inches(4)
        height = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "[Insert Screenshot Here]"
        p.font.size = Pt(24)

# 1. Introduction
add_slide(prs, "Reproducing Aspergillus fumigatus Analysis", [
    "Objective: Reproducing the computational analysis of the paper jof-09-01104-v2.pdf",
    "Focus:",
    "  Variant calling",
    "  Identifying specific mutations related to Azole resistance",
    "Speaker: Scott Cain"
])

# 2. Platforms Overview
add_slide(prs, "Tools and Platforms", [
    "BRC-Analytics",
    "  A platform for organism-centric analysis",
    "  Seamless integration with Galaxy",
    "Galaxy Platform",
    "  Reproducible bioinformatics workflow environment"
])

# 3. Setting Up Analysis in BRC-Analytics
add_slide(prs, "Setting Up Analysis (BRC-Analytics)", [
    "Navigate to brc-analytics.org",
    "Locate Organism:",
    "  Filter and search for 'Aspergillus fumigatus'",
    "Configure Analysis:",
    "  Select 'Aspergillus fumigatus' -> 'Analyze'",
    "  Choose 'Variant calling'",
    "  Enter SRA/ENA Accessions from two BioProjects",
    "  Select all 22 corresponding sequencing runs",
    "  Click 'Launch In Galaxy'"
], screenshot_placeholder=True)

# 4. Running the Haploid Variant Caller in Galaxy
add_slide(prs, "Variant Calling Workflow (Galaxy)", [
    "Workflow Initialization:",
    "  Open pre-configured workflow: 'Paired End Reads'",
    "  Select the collection containing the 22 paired-end datasets",
    "  Execute 'Run Workflow'",
    "Standard Steps Included:",
    "  Quality Control (fastp)",
    "  Mapping",
    "  Variant calling"
])

# 5. Extracting and Processing cyp51A Reads
add_slide(prs, "Extracting cyp51A Reads", [
    "Slicing BAM Files:",
    "  Use 'BAM by genomic regions' tool",
    "  Input coordinates for cyp51A (NC_007197.1 / NC_007196.1)",
    "  Rename output to 'cyp51A-associated reads'",
    "Conversion and Assembly:",
    "  Convert BAM to FASTQ using 'bedtools'",
    "  Nest collection & run 'SPAdes' genome assembler",
    "Finding Tandem Repeats:",
    "  Run 'etandem' on SPAdes scaffolds"
], screenshot_placeholder=True)

# 6. AI-Assisted Scripting: cyp51A Variants GFF3
add_slide(prs, "AI-Assisted Scripting: GFF3 Generation", [
    "Using Gemini CLI:",
    "  Iteratively built a JupyterLite notebook (variants_cyp51A.ipynb)",
    "Notebook Execution:",
    "  Run within Galaxy's JupyterLite interactive environment",
    "  Dynamically extracts variants and tandem repeats",
    "  Maps them back to specific SRA accessions",
    "Output:",
    "  Combined GFF3 file (cyp51A variant script output)"
])

# 7. AI-Assisted Scripting: Finding Resistant-Specific SNPs
add_slide(prs, "AI-Assisted Scripting: Resistant-Specific SNPs", [
    "Identifying Phenotypes:",
    "  Identify accessions for resistant samples (e.g., ERR14230100)",
    "Using Gemini:",
    "  Wrote second JupyterLite script for Set Operations",
    "Set Operations:",
    "  Union of all SNPs in susceptible samples",
    "  Intersection of SNPs across resistant samples",
    "  Calculate set difference for resistant-specific mutations",
    "Output:",
    "  Resistant Specific SNPs (VCF)"
], screenshot_placeholder=True)

# 8. Conclusion
add_slide(prs, "Conclusion", [
    "Summary of Findings:",
    "  Identified mutations correlating with the original paper",
    "The Power of AI + Galaxy:",
    "  Integrating AI tools (Gemini) with platforms (BRC-Analytics/Galaxy)",
    "  Drastically accelerates reproducible research",
    "  Simplifies complex data wrangling and custom scripting"
])

out_path = '/app/presentation_markdown/AsperFest_Presentation.pptx'
prs.save(out_path)
print(f"Presentation saved to {out_path}")
